"""Tests de la API de ToMarkdown (endpoint de conversión)."""

import io
from pathlib import Path

import pytest
from docx import Document
from httpx import AsyncClient
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen.canvas import Canvas

FIXTURES_DIR = Path(__file__).parent / "fixtures"
DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
MAX_FILE_SIZE = 10 * 1024 * 1024


def _make_sample_docx() -> bytes:
    """Genera un DOCX válido en memoria con un encabezado y un párrafo."""
    buffer = io.BytesIO()
    doc = Document()
    doc.add_heading("Titulo de Prueba", level=1)
    doc.add_paragraph("Parrafo de ejemplo.")
    doc.save(buffer)
    return buffer.getvalue()


def _make_sample_pptx() -> bytes:
    """Genera un PPTX válido en memoria con un título y una viñeta."""
    buffer = io.BytesIO()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Título y contenido
    slide.shapes.title.text_frame.text = "Diapositiva Uno"
    slide.placeholders[1].text_frame.text = "Punto de ejemplo"
    prs.save(buffer)
    return buffer.getvalue()


def _make_sample_pdf_with_table() -> bytes:
    """Genera un PDF en memoria con una tabla dibujada a mano (canvas).

    reportlab.Table no sirve como fixture: su GRID no genera líneas
    vectoriales detectables por pymupdf.find_tables(), y la conversión
    quedaría como texto plano. Dibujando las celdas con canvas, PyMuPDF
    detecta la tabla y pymupdf4llm la convierte a Markdown GFM.
    """
    buffer = io.BytesIO()
    c = Canvas(buffer, pagesize=letter)

    # Celdas 3x2: líneas verticales (columnas) y horizontales (filas).
    col_x = [50, 120, 180, 250]
    row_y = [760, 720, 680]
    for x in col_x:
        c.line(x, row_y[0], x, row_y[-1])
    for y in row_y:
        c.line(col_x[0], y, col_x[-1], y)

    c.drawString(55, 740, "Nombre")
    c.drawString(125, 740, "Edad")
    c.drawString(185, 740, "Ciudad")
    c.drawString(55, 700, "Matias")
    c.drawString(125, 700, "30")
    c.drawString(185, 700, "Buenos Aires")

    c.showPage()
    c.save()
    return buffer.getvalue()


def _make_docx_with_table() -> bytes:
    """DOCX con párrafo, tabla y párrafo para validar el orden intercalado."""
    buffer = io.BytesIO()
    doc = Document()
    doc.add_heading("Encabezado", level=1)
    doc.add_paragraph("Parrafo antes de la tabla.")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Nombre"
    table.cell(0, 1).text = "Edad"
    table.cell(1, 0).text = "Ana"
    table.cell(1, 1).text = "25"
    doc.add_paragraph("Parrafo despues de la tabla.")
    doc.save(buffer)
    return buffer.getvalue()


def _make_docx_with_soft_break() -> bytes:
    """DOCX con un soft break (<w:br/>) dentro de un párrafo."""
    buffer = io.BytesIO()
    doc = Document()
    p = doc.add_paragraph()
    p.add_run("Hola")
    p.add_run().add_break()
    p.add_run("Mundo")
    doc.save(buffer)
    return buffer.getvalue()


@pytest.mark.anyio
async def test_health_check(client: AsyncClient) -> None:
    response = await client.get("/")

    assert response.status_code == 200
    assert response.json()["status"] == "online"


@pytest.mark.anyio
async def test_convert_docx_happy_path(client: AsyncClient) -> None:
    response = await client.post(
        "/api/v1/convert",
        files={"file": ("ejemplo.docx", _make_sample_docx(), DOCX_MIME)},
    )

    assert response.status_code == 200
    data = response.json()
    assert "# Titulo de Prueba" in data["markdown"]
    assert data["token_count"] > 0


@pytest.mark.anyio
async def test_convert_pdf_happy_path(client: AsyncClient) -> None:
    pdf_bytes = (FIXTURES_DIR / "sample.pdf").read_bytes()
    response = await client.post(
        "/api/v1/convert",
        files={"file": ("sample.pdf", pdf_bytes, "application/pdf")},
    )

    assert response.status_code == 200
    data = response.json()
    assert "Hola ToMarkdown" in data["markdown"]
    assert data["token_count"] > 0


@pytest.mark.anyio
async def test_convert_pdf_preserves_table(client: AsyncClient) -> None:
    pdf_bytes = _make_sample_pdf_with_table()
    response = await client.post(
        "/api/v1/convert",
        files={"file": ("tabla.pdf", pdf_bytes, "application/pdf")},
    )

    assert response.status_code == 200
    data = response.json()
    md = data["markdown"]

    table_lines = [ln for ln in md.splitlines() if ln.strip().startswith("|")]
    assert table_lines, "la tabla no se convirtió a GFM"
    assert "Nombre" in table_lines[0] and "Ciudad" in table_lines[0]
    assert "Matias" in table_lines[-1] and "Buenos Aires" in table_lines[-1]
    assert any("---" in ln for ln in table_lines)
    assert data["token_count"] > 0


@pytest.mark.anyio
async def test_convert_docx_interleaves_table_and_paragraphs(
    client: AsyncClient,
) -> None:
    docx_bytes = _make_docx_with_table()
    response = await client.post(
        "/api/v1/convert",
        files={"file": ("estructura.docx", docx_bytes, DOCX_MIME)},
    )

    assert response.status_code == 200
    md = response.json()["markdown"]

    # El orden del documento se respeta: texto, tabla GFM, texto.
    antes = md.index("Parrafo antes")
    tabla = md.index("| Nombre |")
    despues = md.index("Parrafo despues")
    assert antes < tabla < despues
    assert "| Nombre | Edad |" in md
    assert "| --- | --- |" in md
    assert "| Ana | 25 |" in md


@pytest.mark.anyio
async def test_convert_docx_preserves_soft_break(client: AsyncClient) -> None:
    docx_bytes = _make_docx_with_soft_break()
    response = await client.post(
        "/api/v1/convert",
        files={"file": ("saltos.docx", docx_bytes, DOCX_MIME)},
    )

    assert response.status_code == 200
    md = response.json()["markdown"]
    # El soft break ya no pega el texto contiguo ("Hola" + "Mundo").
    assert "HolaMundo" not in md
    assert "Hola" in md
    assert "Mundo" in md


@pytest.mark.anyio
async def test_convert_unsupported_format(client: AsyncClient) -> None:
    response = await client.post(
        "/api/v1/convert",
        files={"file": ("archivo.rtf", b"hola", "text/rtf")},
    )

    assert response.status_code == 422


@pytest.mark.anyio
async def test_convert_rejects_file_over_10mb(client: AsyncClient) -> None:
    oversized = b"x" * (MAX_FILE_SIZE + 1)  # 10 MB + 1 byte
    response = await client.post(
        "/api/v1/convert",
        files={"file": ("grande.pdf", oversized, "application/pdf")},
    )

    assert response.status_code == 413


@pytest.mark.anyio
async def test_convert_pptx_happy_path(client: AsyncClient) -> None:
    response = await client.post(
        "/api/v1/convert",
        files={"file": ("presentacion.pptx", _make_sample_pptx(), PPTX_MIME)},
    )

    assert response.status_code == 200
    data = response.json()
    assert "# Diapositiva Uno" in data["markdown"]
    assert "- Punto de ejemplo" in data["markdown"]
    assert data["token_count"] > 0


@pytest.mark.anyio
async def test_convert_txt_happy_path(client: AsyncClient) -> None:
    response = await client.post(
        "/api/v1/convert",
        files={"file": ("notas.txt", b"Hola desde txt", "text/plain")},
    )

    assert response.status_code == 200
    data = response.json()
    assert "Hola desde txt" in data["markdown"]
    assert data["token_count"] > 0


@pytest.mark.anyio
async def test_convert_txt_latin1_encoding(client: AsyncClient) -> None:
    # "Café" en Latin-1 no es UTF-8 válido → el parser hace fallback a Latin-1.
    latin1_bytes = "Café".encode("latin-1")
    response = await client.post(
        "/api/v1/convert",
        files={"file": ("cafe.txt", latin1_bytes, "text/plain")},
    )

    assert response.status_code == 200
    data = response.json()
    assert "Café" in data["markdown"]
    assert data["token_count"] > 0


@pytest.mark.anyio
async def test_rate_limit_returns_429(client: AsyncClient) -> None:
    pdf_bytes = (FIXTURES_DIR / "sample.pdf").read_bytes()
    files = {"file": ("sample.pdf", pdf_bytes, "application/pdf")}

    # Las 10 primeras conversiones del minuto pasan…
    for _ in range(10):
        response = await client.post("/api/v1/convert", files=files)
        assert response.status_code == 200

    # …la 11ª supera el límite (10/minute) y responde 429.
    response = await client.post("/api/v1/convert", files=files)
    assert response.status_code == 429
