"""Parser de presentaciones PowerPoint (.pptx) a Markdown."""

import io

from pptx import Presentation


def parse_pptx_to_markdown(file_stream: io.BytesIO) -> str:
    """Extrae el texto de cada diapositiva y lo estructura en Markdown.

    Cada diapositiva se marca con un H2, el título se renderiza como H1 y
    el resto del contenido como viñetas.
    """
    presentation = Presentation(file_stream)
    markdown_lines = []

    for idx, slide in enumerate(presentation.slides, start=1):
        title = None
        bullets: list[str] = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text = shape.text_frame.text.strip()
            if not text:
                continue

            # El placeholder con índice 0 es el título de la diapositiva.
            if (
                shape.is_placeholder
                and shape.placeholder_format.idx == 0
                and title is None
            ):
                title = text
            else:
                bullets.extend(
                    line.strip() for line in text.splitlines() if line.strip()
                )

        markdown_lines.append(f"## Diapositiva {idx}\n")
        if title:
            markdown_lines.append(f"# {title}\n")
        for bullet in bullets:
            markdown_lines.append(f"- {bullet}\n")

    return "\n".join(markdown_lines)
